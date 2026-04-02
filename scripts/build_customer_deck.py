"""
Generate VitalBee customer presentation (PowerPoint).
Run from repo root: python scripts/build_customer_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1]
    tf.text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets, subtitle=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
    if subtitle:
        # add a small footer note as extra paragraph
        p = tf.add_paragraph()
        p.text = subtitle
        p.level = 0
        p.font.size = Pt(14)
        p.font.italic = True
    return slide


def add_two_column_bullets(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Title + content layout with manual text boxes for two columns."""
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True

    def fill_column(x, y, w, h, heading, items):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading
        p.font.size = Pt(22)
        p.font.bold = True
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(17)
            p.space_before = Pt(6)

    fill_column(0.5, 1.35, 4.4, 5.5, left_title, left_bullets)
    fill_column(5.1, 1.35, 4.4, 5.5, right_title, right_bullets)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "VitalBee",
        "Drone-based survivor detection for Search & Rescue\nTeam DİH · Istanbul Technical University (ITU)\nvitalbee@itu.edu.tr",
    )

    add_bullet_slide(
        prs,
        "The challenge",
        [
            "After earthquakes, collapses, and landslides, every minute counts.",
            "Manual search is dangerous, slow, and hard to scale across wide debris fields.",
            "Teams often juggle separate tools: aerial drones vs. ground sensors—with fragmented data and heavy operator load.",
        ],
    )

    add_bullet_slide(
        prs,
        "What VitalBee delivers",
        [
            "An intelligent aerial platform that scans wide areas autonomously.",
            "Signal-based sensing plus AI-assisted fusion—not just raw video or thermal feeds.",
            "Live geo-referenced probability maps, ranked priority points, and short evidence summaries for dispatch decisions.",
        ],
    )

    add_two_column_bullets(
        prs,
        "How the system works (overview)",
        "Wide-area phase",
        [
            "Adaptive grid-based search pattern",
            "RF, acoustic, optional thermal sensing",
            "Onboard edge compute (e.g. Jetson-class)",
            "Preliminary geo-referenced probability map",
        ],
        "Targeted phase",
        [
            "Closer verification over high-interest cells",
            "Controlled probing and temporal analysis",
            "Reduced false positives via triangulation",
            "Confidence-ranked outputs to the field dashboard",
        ],
    )

    add_bullet_slide(
        prs,
        "Why teams would choose VitalBee",
        [
            "End-to-end workflow: detection → localization → prioritization → mission logging.",
            "Modular payloads and models—extend without rebuilding core mapping software.",
            "Built for low-visibility, structurally unsafe environments where speed and clarity matter.",
        ],
    )

    add_bullet_slide(
        prs,
        "Who we serve",
        [
            "National and municipal disaster response agencies",
            "Civil protection and public-safety UAV programs",
            "Humanitarian organizations needing faster, coordinated SAR",
        ],
        "Go-to-market: government-backed R&D, pilots, and institutional procurement channels.",
    )

    add_bullet_slide(
        prs,
        "Commercial model (illustrative)",
        [
            "Drone kit (hardware): $15,000 per unit — one-time",
            "Software platform: $3,000 annual subscription per unit",
            "Optional maintenance & training: $2,000 per client per year",
            "Hybrid model: hardware sale + recurring software and services for predictable long-term revenue",
        ],
        "Figures reflect early business-model coursework assumptions; final pricing subject to pilot feedback and certification scope.",
    )

    add_bullet_slide(
        prs,
        "Delivery & trust",
        [
            "Onboarding, training, technical support, and continuous software updates",
            "Phased prototyping, safety gates, lab and debris-site testing",
            "Explainable scoring so operators can trust recommendations under pressure",
        ],
    )

    add_bullet_slide(
        prs,
        "Strategic collaboration (roadmap)",
        [
            "Sensors & rugged hardware: STM Defence, Aselsan",
            "Edge AI & acceleration: NVIDIA",
            "Reliable, explainable AI: Anthropic",
            "Universities and public institutions: validation pilots and co-funded trials",
        ],
    )

    add_bullet_slide(
        prs,
        "Risks we acknowledge",
        [
            "Strong competition from established SAR and drone vendors",
            "Public procurement can be slow—cash flow and pipeline planning matter",
            "Field validation is essential before mission-critical adoption",
        ],
        "Mitigations: subscription/leasing or drone-as-a-service options, tiered analytics, and demonstration-led sales.",
    )

    add_bullet_slide(
        prs,
        "Next steps with your organization",
        [
            "Introduce use cases and operating environments (urban collapse, landslide, etc.)",
            "Define pilot scope, safety requirements, and data-sharing needs",
            "Align on procurement path and timeline for a controlled field trial",
        ],
    )

    # Closing / contact
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Thank you"
    slide.placeholders[1].text = "VitalBee · Team DİH · ITU\nvitalbee@itu.edu.tr\n\nWe welcome conversations on pilots, partnerships, and procurement."

    out = Path(__file__).resolve().parent.parent / "VitalBee-customer-presentation.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
